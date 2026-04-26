#!/usr/bin/env python3
"""
Build updated DAC Presentation with corrected figures and new slides.
Replaces old figures, adds experimental design and kinetics slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

FIGDIR = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/figures'
PPTX_IN = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/DAC_Presentation.pptx'
PPTX_OUT = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/DAC_Presentation_v2.pptx'

prs = Presentation(PPTX_IN)
slide_width = prs.slide_width
slide_height = prs.slide_height

# --- Helper functions ---
def find_and_replace_image(slide, new_image_path):
    """Replace the first image found on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(new_image_path, left, top, width, height)
            return True
    return False

def add_slide_with_title_and_content(prs, title_text, content_lines, layout_idx=1):
    """Add a new slide with title and bullet content."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
        elif shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(content_lines):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = line
                para.font.size = Pt(16)
                if line.startswith('•') or line.startswith('-'):
                    para.level = 1
                    para.font.size = Pt(14)
    return slide

def add_slide_with_image(prs, title_text, image_path, subtitle_text=None):
    """Add a new slide with title and a centered image."""
    layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(layout)

    # Title
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = title_text
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    para.alignment = PP_ALIGN.CENTER

    # Image
    img_top = Inches(1.1)
    img_max_w = Inches(11.5)
    img_max_h = Inches(5.8)

    from PIL import Image
    with Image.open(image_path) as img:
        iw, ih = img.size
    aspect = iw / ih

    if aspect > (11.5 / 5.8):
        w = img_max_w
        h = int(w / aspect)
    else:
        h = img_max_h
        w = int(h * aspect)

    left = int((slide_width - w) / 2)
    slide.shapes.add_picture(image_path, left, img_top, w, h)

    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(10)
        sp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        sp.font.italic = True
        sp.alignment = PP_ALIGN.CENTER

    return slide


# === STEP 1: Replace figures on existing slides ===
figure_map = {
    6: 'kinetic_proofreading.png',
    7: 'optimal_dwell_time.png',
    8: 'serial_engagement.png',
    10: 'catch_vs_slip_bond.png',
    12: 'affinity_windows.png',
    18: 'experimental_workflow.png',
}

for slide_num, fig_name in figure_map.items():
    slide = prs.slides[slide_num - 1]
    fig_path = os.path.join(FIGDIR, fig_name)
    if os.path.exists(fig_path):
        if find_and_replace_image(slide, fig_path):
            print(f'  Replaced figure on slide {slide_num}: {fig_name}')
        else:
            print(f'  WARNING: No image found on slide {slide_num} to replace')

# Check for car_vs_tcr_affinity on slide 15
slide15 = prs.slides[14]
car_fig = os.path.join(FIGDIR, 'car_vs_tcr_affinity.png')
if os.path.exists(car_fig):
    if find_and_replace_image(slide15, car_fig):
        print(f'  Replaced figure on slide 15: car_vs_tcr_affinity.png')

# Also replace signaling_cascade on slide 11 if present
slide11 = prs.slides[10]
sig_fig = os.path.join(FIGDIR, 'signaling_cascade.png')
if os.path.exists(sig_fig):
    if find_and_replace_image(slide11, sig_fig):
        print(f'  Replaced figure on slide 11: signaling_cascade.png')

# Replace tcr_synapse on slide 9 if present
slide9 = prs.slides[8]
syn_fig = os.path.join(FIGDIR, 'tcr_synapse_organization.png')
if os.path.exists(syn_fig):
    if find_and_replace_image(slide9, syn_fig):
        print(f'  Replaced figure on slide 9: tcr_synapse_organization.png')


# === STEP 2: Remove old "Thank You" slide (last slide) ===
# We'll add new slides and a new Thank You at the end
last_slide = prs.slides[len(prs.slides) - 1]
rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')


# === STEP 3: Add new slides BEFORE Thank You ===
# We need to insert before the last slide. In python-pptx, slides are appended.
# Strategy: remove Thank You, add new slides, re-add Thank You.

# Remove last slide (Thank You)
prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
prs.part.drop_rel(rId)
print('\n  Removed old Thank You slide')

# --- New Slide: NNK Screening Strategy ---
add_slide_with_image(prs, 'Three-Phase NNK Library Screening Strategy',
                     os.path.join(FIGDIR, 'screening_strategy.png'),
                     'Coverage: Pines et al., 2022, Synth Biol (PMC9205323)')
print('  Added: NNK Screening Strategy')

# --- New Slide: FACS Panel Design ---
add_slide_with_image(prs, 'FACS Panel Design for CAR-T Characterization',
                     os.path.join(FIGDIR, 'facs_panel_summary.png'),
                     'All clones verified via BioLegend, Miltenyi, Thermo Fisher product pages')
print('  Added: FACS Panel Design')

# --- New Slide: Published FMC63 Variant Affinities ---
add_slide_with_image(prs, 'Published FMC63 Affinity Variants at Target Residues',
                     os.path.join(FIGDIR, 'fmc63_variant_affinities.png'),
                     'Data: Singh et al., 2023, Science Immunology (PMC10228544); '
                     'Seigner et al., 2023, Sci Rep (PMC10754921)')
print('  Added: FMC63 Variant Affinities')

# --- New Slide: Rechallenge & Controls ---
slide_rc = add_slide_with_title_and_content(prs,
    'Rechallenge Assay & Essential Controls',
    [
        'Rechallenge Protocol (Wang et al., 2019, J ImmunoTher Cancer):',
        '• E:T 1:4 → fresh Raji added every 2 days → 4 rounds over 7 days',
        '• Readouts: activation (CD69, 4-1BB), exhaustion (PD-1/TIM-3/LAG-3), memory (CD45RO/CD62L)',
        '',
        'Essential Controls (8 conditions):',
        '• Un-transduced Jurkat + Raji (non-specific activation baseline)',
        '• Mock-transduced (empty vector) + Raji (transduction effect)',
        '• CAR-Jurkat + K562 (CD19-negative specificity control)',
        '• CAR-Jurkat alone (tonic signaling assessment)',
        '• Raji alone (spontaneous death)',
        '• WT FMC63 CAR + Raji (positive benchmark)',
        '• PMA/ionomycin-treated Jurkat (maximum activation)',
        '',
        'Statistics: One-way ANOVA + Dunnett\'s post-hoc (many-to-one vs WT)',
        '• n=3 independent transductions; triplicate wells',
    ])
print('  Added: Rechallenge & Controls')

# --- New Slide: Antigen Density ---
add_slide_with_image(prs, 'Antigen Density Determines CAR-T Activation',
                     os.path.join(FIGDIR, 'antigen_density_threshold.png'),
                     'Adapted from: Majzner et al., 2020, Cancer Discovery (PMC7939454)')
print('  Added: Antigen Density')

# --- New Slide: Alternative Systems ---
slide_alt = add_slide_with_title_and_content(prs,
    'Alternative Systems for Comprehensive Validation',
    [
        'Raji limitations: fixed CD19 density, FACS quantification challenges',
        '',
        'Proposed complementary systems:',
        '• Raji-GFP-Luc2 (ATCC CCL-86-GFP-LUC2) — immediate: GFP for FACS + luciferase for killing',
        '• NALM-6 CD19-KO + graded re-expression — tunable CD19 density (45-45,851 mol/cell)',
        '  (Majzner et al., 2020, Cancer Discovery 10:702)',
        '• CHO-CD19 Low/Med/High — adherent targets for real-time xCELLigence monitoring',
        '',
        'FACS solutions:',
        '• Counting beads (BioLegend Precision Count Beads) for absolute quantification',
        '• CellTrace Violet pre-labeling of targets for surface marker-independent gating',
        '',
        'Advanced (future):',
        '• Supported lipid bilayers + TIRF for CAR synapse imaging',
        '• 2D kinetics (micropipette/BFP) — no published data for ANY CAR system (novelty)',
    ])
print('  Added: Alternative Systems')

# --- New Slide: Kinetics Platforms ---
add_slide_with_image(prs, 'Binding Kinetics: SPR vs BLI Platform Comparison',
                     os.path.join(FIGDIR, 'platform_comparison.png'),
                     'FMC63-CD19 KD = 5.1 nM (Seigner et al., 2023, Biacore T200)')
print('  Added: Kinetics Platforms')

# --- New Slide: Expected Outcomes (expanded) ---
slide_out = add_slide_with_title_and_content(prs,
    'Expected Outcomes & Significance',
    [
        'Expected Outcomes:',
        '• Library of anti-CD19 CAR variants with systematically altered binding kinetics',
        '• Comprehensive functional characterization (activation, exhaustion, memory, serial killing)',
        '• Quantitative kinetic data (KD, kon, koff) via SPR/BLI for selected variants',
        '• Identification of the kinetic parameter(s) that best predict CAR-T functional outcome',
        '• Definition of the optimal affinity window for anti-CD19 CAR-T therapy',
        '',
        'Significance:',
        '• Bridges fundamental TCR-pMHC biology with translational CAR-T design',
        '• Published data validates residue choice: Y260A/Y261A (Singh et al., 2023)',
        '• Rational framework for affinity-optimized CAR engineering',
        '• Generalizable approach applicable to CARs targeting other antigens',
        '• Potential to improve CAR-T persistence, reduce exhaustion & toxicity',
    ])
print('  Added: Expected Outcomes')

# --- New Slide: Thank You ---
layout = prs.slide_layouts[5]  # Blank
slide_ty = prs.slides.add_slide(layout)
txBox = slide_ty.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
para = tf.paragraphs[0]
para.text = 'Thank You'
para.font.size = Pt(44)
para.font.bold = True
para.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
para.alignment = PP_ALIGN.CENTER

sub = slide_ty.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1.5))
stf = sub.text_frame
stf.word_wrap = True
lines = [
    'PI: Dr. Kausik Chakraborty | Co-PI: Dr. Ankesh Kumar Jaiswal',
    'DAC Members: Dr. Arpan Parichha | Dr. Chetana Sachidanandan | Dr. Sheetal Gandotra',
    '',
    'Manpreet Kour | CSIR-IGIB | AcSIR Reg. 10BB25J02028'
]
for i, line in enumerate(lines):
    if i == 0:
        p = stf.paragraphs[0]
    else:
        p = stf.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p.alignment = PP_ALIGN.CENTER
print('  Added: Thank You')

# === STEP 4: Save ===
prs.save(PPTX_OUT)
print(f'\nSaved: {PPTX_OUT}')
print(f'Total slides: {len(prs.slides)}')
