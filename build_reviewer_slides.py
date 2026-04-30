#!/usr/bin/env python3
"""
Build integration-ready PPT for DAC: Experimental Design — Reviewer Considerations.
Designed to be inserted into Manpreet's existing DAC_Presentation_v2.pptx.

Slides covered:
  1. Title
  2. Bottom-line summary (5 critical issues)
  3. Revised experimental workflow
  4. Mouse model choice (NSG vs NSG-MHC-DKO)
  5. Activation / exhaustion kinetic readout (day 1, 2, 4)
  6. FACS panel for activation / exhaustion
  7. Memory phenotyping + variant NGS strategy
  8. Pooled-screen NGS readout (barcode + UMI + MAGeCK)
  9. CD19 production: Expi293F + SF-CD19 (Laurent 2021)
 10. SPR protocol (Seigner 2023 parameters)
 11. Tiered biophysics: BLI -> SPR -> 2D -> BFP
 12. Citation correction: He et al. 2023 (NOT Singh)
 13. Summary recommendations table
 14. References
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/Experimental_Design_Reviewer_Slides.pptx'

# Color palette
NAVY = RGBColor(0x1A, 0x23, 0x7E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_LIGHT = RGBColor(0xE8, 0xEE, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_title(slide, text, top=Inches(0.3), height=Inches(0.7), size=28, color=NAVY, align=PP_ALIGN.LEFT, left=Inches(0.5), width=None):
    if width is None:
        width = SW - Inches(1.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    return box


def add_subtitle(slide, text, top, size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.LEFT, left=Inches(0.5), width=None):
    if width is None:
        width = SW - Inches(1.0)
    box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    return box


def add_bullets(slide, lines, top, left=Inches(0.5), width=None, height=None, size=14, indent_size=12):
    if width is None:
        width = SW - Inches(1.0)
    if height is None:
        height = SH - top - Inches(0.4)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Bullet character based on level
        if level == 0:
            p.text = u'•  ' + text
            font_size = size
        elif level == 1:
            p.text = u'    ◦  ' + text
            font_size = indent_size
        else:
            p.text = u'        –  ' + text
            font_size = indent_size - 1
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK_GRAY
            r.font.name = 'Calibri'
        p.space_after = Pt(4)
    return box


def add_table(slide, data, top, left, width, height, header_color=NAVY, alt_row=BLUE_LIGHT, font_size=11, header_white=True):
    """Add a table. data is list of lists; first row is header."""
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.name = 'Calibri'
                    if ri == 0:
                        r.font.bold = True
                        if header_white:
                            r.font.color.rgb = WHITE
                    else:
                        r.font.color.rgb = DARK_GRAY
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row
    return tbl_shape


def add_callout_box(slide, text, top, left, width, height, fill_color, border_color=None, font_color=WHITE, font_size=14, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = border_color
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = font_color
        r.font.name = 'Calibri'
    return rect


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), SH - Inches(0.35), SW - Inches(1.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(9)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
        r.font.name = 'Calibri'
    return box


# =====================================================================
# Slide 1: Title
# =====================================================================
slide = prs.slides.add_slide(BLANK)

# Background accent stripe
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(2.5))
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), SW - Inches(1.4), Inches(1.6))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Experimental Design — Reviewer Considerations'
for r in p.runs:
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = 'Calibri'

p2 = tf.add_paragraph()
p2.text = 'Confounder audit and refinements for the in vivo + in vitro + biophysics arms'
for r in p2.runs:
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = WHITE
    r.font.name = 'Calibri'

# Subtitle/info block
info = slide.shapes.add_textbox(Inches(0.7), Inches(2.9), SW - Inches(1.4), Inches(2.5))
tf = info.text_frame
tf.word_wrap = True
lines = [
    ('Manpreet Kour | PhD Scholar | CSIR-IGIB | AcSIR Reg. 10BB25J02028', True, 18, NAVY),
    ('PI: Dr. Kausik Chakraborty (Chief Scientist, CSIR-IGIB)', False, 14, DARK_GRAY),
    ('Co-PI: Dr. Ankesh Kumar Jaiswal (Project Scientist)', False, 14, DARK_GRAY),
    ('', False, 12, DARK_GRAY),
    ('System: Anti-CD19 CAR (FMC63 scFv)  |  Target: CD19 (PDB 7URV)', False, 14, DARK_GRAY),
    ('NNK saturation library: S214, Trp212 (primary) | Y260, Y261 (secondary)', False, 14, DARK_GRAY),
    ('', False, 12, DARK_GRAY),
    ('Purpose of these slides: integration-ready content for DAC presentation', False, 13, DARK_GRAY),
    ('All citations directly verified against PubMed', False, 13, GREEN_ACCENT),
]
for i, (text, bold, size, color) in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    p.space_after = Pt(4)

add_footer(slide, '2026-05-01  |  Document 18 + reviewer slides  |  visvikbharti/DAC-Meeting-CAR-T')

# =====================================================================
# Slide 2: Bottom-line summary — 5 critical issues
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Five issues to address before DAC presentation', size=26)
add_subtitle(slide, 'Top-priority refinements identified by independent expert review of the proposed design', Inches(1.0))

issues = [
    ('1. Citation correction', 'PDB 7URV / Sci Immunol 2023 paper has first author HE (not Singh). Repo + slide deck use "Singh et al." in 16 files — must be corrected.', RED_ACCENT),
    ('2. Sampling timing (clarified)', 'Day 1, 2, 4 = activation / exhaustion kinetics (correct). Memory phenotype = terminal sacrifice (end of 21-d rest, post-rechallenge). Do NOT sort "memory subsets" at day 1-4 post-rechallenge — those are recall effectors.', ORANGE),
    ('3. "Humanized mouse" terminology', 'For adoptive CAR-T transfer to NALM-6, use NSG (or NSG-MHC-DKO if rest >4 wk; Brehm 2019, PMID 30383447). PBMC- or HSC-humanized mice will develop xeno-GvHD that masks the readout.', ORANGE),
    ('4. Pre-injection library NGS is mandatory', 'Without sequencing the input plasmid library + post-expansion infusion product, no enrichment ratio is interpretable. Currently not in the plan.', RED_ACCENT),
    ('5. CD19 production: HEK293S vs Expi293F', '"HEK293S" in the field = GnTI⁻ line (high-mannose only, used for crystallography). For SPR/BLI kinetics use Expi293F + SF-CD19 stabilized monomer (Laurent 2021, PMID 33843201).', RED_ACCENT),
]

y = Inches(1.6)
for title_text, body, color in issues:
    # Color stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.15), Inches(1.0))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    # Title + body
    box = slide.shapes.add_textbox(Inches(0.8), y, SW - Inches(1.3), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_text
    for r in p1.runs:
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = body
    for r in p2.runs:
        r.font.size = Pt(12)
        r.font.color.rgb = DARK_GRAY
        r.font.name = 'Calibri'
    y += Inches(1.05)

add_footer(slide, 'Source: 18_Experimental_Design_Expert_Review.md §1, §8')

# =====================================================================
# Slide 3: Revised experimental workflow
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Revised experimental workflow (recommended)')
add_subtitle(slide, 'Each timepoint matched to its biologically appropriate readout', Inches(1.0))

flow_data = [
    ['Phase', 'Timing', 'Sample', 'Readout', 'Question answered'],
    ['Library QC', 'Pre-injection', 'Plasmid + infusion product', 'Amplicon / barcode NGS', 'Variant representation in input (REQUIRED reference)'],
    ['Primary kinetics', 'Day 1, 2, 4 post-CAR-T', 'Peripheral blood (serial)', 'FACS panel (activation + exhaustion); sort + NGS', 'Activation / exhaustion KINETICS per variant'],
    ['Tissue infiltration', 'Day 4, 14 (terminal cohort)', 'Spleen + bone marrow', 'Same panel + counts', 'Tissue distribution of phenotype per variant'],
    ['Memory pool', 'End of day-21 rest (terminal cohort)', 'Spleen + LN + BM', 'Memory subset sort (Tn / Tscm / Tcm / Tem) + NGS', 'Which variants ENRICH in resting memory'],
    ['Recall response', 'Day 1, 2, 4 post-rechallenge', 'Blood (serial)', 'Activation + Ki67 + cytokines', 'Recall response kinetics per variant'],
    ['Recall functional', 'Day 4-7 post-rechallenge (terminal)', 'Spleen + BM + tumor sites', 'Total CAR+ + CAR+Ki67+', 'Which variants drive strongest recall'],
]
add_table(slide, flow_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(4.6), font_size=11)

# Key change box
add_callout_box(slide, 'KEY CHANGE: separate the kinetic question (day 1-4) from the memory-pool question (end-of-rest sacrifice). Different questions, different timepoints.',
                Inches(6.2), Inches(0.5), SW - Inches(1.0), Inches(0.7),
                fill_color=GREEN_ACCENT, font_size=13)

add_footer(slide, 'Source: doc 18 §3a.6, §3.2')

# =====================================================================
# Slide 4: Mouse model choice
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Mouse model choice — clarify "humanized mouse"')
add_subtitle(slide, 'For adoptive CAR-T transfer to NALM-6, full immune reconstitution is NOT needed and is usually counterproductive', Inches(1.0))

mouse_data = [
    ['Strain', 'Use case', 'xeno-GvHD risk during 21-d rest', 'Verdict for this design', 'Citation'],
    ['NSG', 'Standard adoptive CAR-T xenograft', 'Minimal if only CAR-T transferred', 'PRIMARY CHOICE for ≤21-d rest', 'Shultz 2012 PMID 23059428'],
    ['NSG-MHC-DKO', 'Long-window CAR-T xenograft', 'Substantially delayed (>10 wk possible)', 'BEST CHOICE if rest >4 wk', 'Brehm 2019 PMID 30383447'],
    ['NSG-SGM3', 'HSC-myeloid studies', 'HLH-like in HSC-engrafted', 'Not needed here', 'Wunderlich 2018 PMID 30586420'],
    ['hu-PBL-NSG (PBMC)', 'Human T cell biology', 'Onset ~28 d → masks memory readout', 'INADEQUATE — avoid', 'King 2009 PMID 19426570'],
    ['BLT', 'HLA-restricted T cell biology', 'BLT-GvHD by 16-24 wk', 'Not needed here', 'Lan 2006 PMID 16778179'],
]
add_table(slide, mouse_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(3.6), font_size=11)

# NALM-6 parameters mini-box
add_callout_box(slide, 'NALM-6: pre-B ALL, CD19+, HLA-A*02:01. Standard IV dose 0.5-1.0 × 10⁶ in NSG (Brentjens 2007 PMID 17855649; Milone 2009 PMID 19384291). Median survival untreated ≈21-25 d.',
                Inches(5.3), Inches(0.5), SW - Inches(1.0), Inches(0.6),
                fill_color=BLUE_LIGHT, font_color=DARK_GRAY, font_size=11, bold=False)

add_callout_box(slide, 'Recommendation: NSG primary; switch to NSG-MHC-DKO if any rest period exceeds 4 weeks. Drop the term "humanized mouse" unless PBMC/HSC engraftment is intentional.',
                Inches(6.1), Inches(0.5), SW - Inches(1.0), Inches(0.7),
                fill_color=GREEN_ACCENT, font_size=13)

add_footer(slide, 'Source: doc 18 §2.1, §2.2')

# =====================================================================
# Slide 5: Activation / exhaustion kinetic readout
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Activation / exhaustion kinetics — day 1, 2, 4 sampling')
add_subtitle(slide, 'Day 1, 2, 4 IS the right window for these markers (not for memory phenotyping)', Inches(1.0))

kin_data = [
    ['Day', 'Activation markers', 'Effector function', 'Early exhaustion'],
    ['Day 1 (~24 h)', 'CD69 peak; CD25 ramping; ICOS / OX40 induction', 'IFN-γ, TNF-α, IL-2 initiating', 'PD-1 mRNA induction; protein begins to rise'],
    ['Day 2 (~48 h)', 'CD25 plateau; CD69 declining; ICOS / OX40 sustained', 'Granzyme B / perforin granule loading', 'PD-1 surface protein elevated; TIM-3 / LAG-3 emerging'],
    ['Day 4', 'Effector differentiation', 'IFN-γ / TNF-α / cytotoxicity peak', 'PD-1, TIM-3, LAG-3 sustained; TOX accumulating'],
]
add_table(slide, kin_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(2.5), font_size=11)

add_callout_box(slide, 'Sort-then-NGS strategy: at each timepoint, sort CAR+ cells into bins (CD69+ vs CD69-, PD-1ʰⁱ vs PD-1ˡᵒ, TOX+ vs TOX-) → variant NGS in each bin → identifies which scFv mutants enrich in each phenotype. Avoids the bulk-FACS averaging problem.',
                Inches(4.2), Inches(0.5), SW - Inches(1.0), Inches(0.9),
                fill_color=NAVY, font_size=12)

# Sampling strategy
add_subtitle(slide, 'Sampling strategy (recommended hybrid):', Inches(5.3), size=14, italic=False, color=NAVY)
add_bullets(slide, [
    'Peripheral blood serial sampling on day 1, 2, 4, 7 (kinetic resolution; non-terminal)',
    'Terminal cohorts on day 4 + day 14 — spleen + bone marrow (intracellular markers + tissue infiltration)',
    'Pair every kinetic timepoint with a parallel terminal cohort if intracellular markers (TOX, NR4A1, granzyme B, IFN-γ) are needed',
], Inches(5.7), size=13)

add_footer(slide, 'Source: doc 18 §3a.1-3a.4. Long et al. 2015 Nat Med PMID 25939063 (4-1BB exhaustion biology, verified)')

# =====================================================================
# Slide 6: FACS panel for activation / exhaustion
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'FACS panel for activation / exhaustion (CAR+ gated)')
add_subtitle(slide, 'Split across 3-4 panels because of fluorochrome limits', Inches(1.0))

panel_data = [
    ['Category', 'Markers', 'Notes'],
    ['Live / lineage', 'LIVE/DEAD, CD45, CD3, CD8, CD4', 'First gate of any CAR-T panel'],
    ['CAR detection', 'CD19-Fc tetramer (preferred), or anti-idiotype, or surrogate marker', 'Confirm what is in Manpreet\'s vector'],
    ['Activation (surface)', 'CD69, CD25, ICOS, OX40 (CD134)', 'CD69 earliest (peak ~24 h); ICOS / OX40 sustained 48-72 h'],
    ['Exhaustion (surface)', 'PD-1, TIM-3, LAG-3, TIGIT, CD39, 2B4', 'Co-expressed in terminal exhaustion'],
    ['Exhaustion (intracellular)', 'TOX, NR4A1', 'Master regulators; require fix / perm'],
    ['Effector (intracellular)', 'IFN-γ, TNF-α, IL-2, granzyme B, perforin', 'Brief stim + protein transport block before fix'],
    ['Differentiation TFs (intracellular)', 'T-bet, EOMES, BLIMP-1, TCF1 / TCF7', 'TCF1 marks stem-like / progenitor exhausted'],
    ['Proliferation', 'Ki67', 'Intracellular'],
    ['Apoptosis', 'Annexin V (or active caspase-3)', 'Surface stain pre-fix'],
]
add_table(slide, panel_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(5.0), font_size=11)

add_callout_box(slide, 'Confounder: CD19 binding internalizes the CAR — staining with CD19-Fc tetramer at saturating concentration mitigates; consider T2A-fluorescent reporter co-expression for orthogonal CAR detection.',
                Inches(6.6), Inches(0.5), SW - Inches(1.0), Inches(0.6),
                fill_color=ORANGE, font_size=12)

add_footer(slide, 'Source: doc 18 §3a.3. Refs: Wherry & Kurachi 2015 (exhaustion review).')

# =====================================================================
# Slide 7: Memory phenotyping + variant NGS
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Memory phenotyping + variant NGS (terminal sacrifice)')
add_subtitle(slide, 'CD95 is the discriminating marker for Tscm vs naïve — must be in the panel', Inches(1.0))

mem_data = [
    ['Subset', 'Phenotype (CAR+ gated)', 'Defining ref'],
    ['Naïve (Tn)', 'CD45RA+ CD45RO- CCR7+ CD62L+ CD27+ CD28+ CD95 low/neg', 'Sallusto 1999 PMID 10537110'],
    ['Tscm (stem cell memory)', 'CD45RA+ CCR7+ CD62L+ CD27+ CD28+ CD95+ IL-7Rα+', 'Gattinoni 2011 PMID 21926977'],
    ['Tcm (central memory)', 'CD45RA- CD45RO+ CCR7+ CD62L+ CD27+', 'Sallusto 1999 PMID 10537110'],
    ['Tem (effector memory)', 'CD45RA- CD45RO+ CCR7- CD62L-', 'Sallusto 1999 PMID 10537110'],
    ['Trm (tissue-resident, BM)', 'CD69+ CD103+/-', 'Mackay 2013 PMID 24162776'],
]
add_table(slide, mem_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(2.7), font_size=11)

add_subtitle(slide, 'Sort + NGS workflow:', Inches(4.4), size=14, italic=False, color=NAVY)
add_bullets(slide, [
    'Sacrifice cohort at end of day-21 rest (BEFORE rechallenge) — primary memory question',
    'Harvest spleen + LN + bone marrow → enzymatic dissociation, viability stain, FACS sort each subset to ≥95% purity',
    'Per-mouse processing for Tcm/Tem (n ≥ 5); pool 2-3 mice for Tscm if cell numbers tight',
    'Genomic DNA → variant amplicon (or barcode) NGS with UMIs → MAGeCK MLE for enrichment',
    'Cell numbers per mouse: ~1-2 × 10⁶ Tscm achievable from spleen alone (Sabatino 2016 PMID 27226436)',
], Inches(4.8), size=12)

add_footer(slide, 'Source: doc 18 §3.1, §3.2, §3.5. Lugli 2013 Nat Protoc — Tscm sorting SOP.')

# =====================================================================
# Slide 8: Pooled NGS readout
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Pooled-screen NGS readout — barcoding vs amplicon')
add_subtitle(slide, 'Critical first step: map S214, Trp212, Y260, Y261 onto FMC63 VL-218-VH sequence to confirm if they fit in a single 300-bp read', Inches(1.0))

ngs_data = [
    ['Feature', 'Amplicon-seq of mutated region', 'DNA barcoding (5\'/3\' UTR)'],
    ['Cloning effort upfront', 'Low', 'High (variant ↔ barcode lookup)'],
    ['Sequencing complexity', 'Single MiSeq if positions in 300 bp', 'Trivial mapping'],
    ['Errors near NNK', 'Confounding (need UMIs)', 'Eliminated'],
    ['Chimeric reads', 'Risk', 'Negligible'],
    ['Long-read needed if positions > 250 bp apart', 'Yes (PacBio / Nanopore)', 'No'],
    ['Recommended for ~376 variants', 'OK if single amplicon', 'PREFERRED'],
]
add_table(slide, ngs_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(3.0), font_size=11)

# Key requirements
add_callout_box(slide, 'NON-NEGOTIABLE: pre-injection library NGS + post-expansion infusion-product NGS as reference. Without them, no enrichment ratio is interpretable.',
                Inches(4.7), Inches(0.5), SW - Inches(1.0), Inches(0.6),
                fill_color=RED_ACCENT, font_size=13)

# Statistics
add_subtitle(slide, 'Depth and statistics:', Inches(5.4), size=14, italic=False, color=NAVY)
add_bullets(slide, [
    '≥1000 reads per variant per condition (minimum for confident log2-fold-change)',
    'UMIs (6-10 nt) on gene-specific PCR primer — Kivioja 2011 PMID 22101854',
    'MAGeCK MLE for enrichment statistics — Li 2014 PMID 25476604',
    'In vivo CRISPR-screen template: Belk 2022 PMID 35750052',
    'n ≥ 5 mice per arm; pre-registered mixed-effects model (variant ~ time + tissue + (1|mouse))',
], Inches(5.8), size=12)

add_footer(slide, 'Source: doc 18 §3.3, §3.4')

# =====================================================================
# Slide 9: CD19 production — Expi293F + SF-CD19
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'CD19 production — Expi293F + SF-CD19 (Laurent 2021)')
add_subtitle(slide, 'Wrong cell line will cost months of aggregation troubleshooting', Inches(1.0))

cd19_data = [
    ['Cell line', 'Glycoform', 'Yield', 'Monomer', 'Verdict'],
    ['Expi293F + SF-CD19', 'Full complex (matches in vivo)', 'High (10-50 mg/L)', 'Excellent (engineered)', 'PRIMARY CHOICE'],
    ['HEK293-6E + SF-CD19', 'Full complex', 'Moderate-high', 'Excellent', 'Equivalent fallback (NRC license)'],
    ['HEK293F (FreeStyle)', 'Full complex', 'Moderate', 'Excellent (with SF-CD19)', 'OK if Expi293F unavailable'],
    ['HEK293S GnTI⁻', 'Man₅ ONLY — wrong glycoform', 'Moderate', '–', 'AVOID (crystallography only)'],
    ['CHO', 'Slightly different (LacdiNAc)', 'High', '–', 'Avoid for kinetics'],
    ['Sf9 / E. coli', 'Wrong (or no) glycans', '–', '–', 'AVOID for CD19'],
]
add_table(slide, cd19_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(3.0), font_size=11)

add_callout_box(slide, 'SF-CD19: SuperFolder stabilized monomer from Laurent E et al. 2021 ACS Synth Biol 10:1184-1198, PMID 33843201. Yeast-display directed evolution → soluble monomeric CD19 (>99% by SEC-MALS). Solves WT-CD19-ECD aggregation. Senior author Traxlmayr — same group as Seigner 2023.',
                Inches(4.7), Inches(0.5), SW - Inches(1.0), Inches(0.85),
                fill_color=GREEN_ACCENT, font_size=12)

# Construct details
add_subtitle(slide, 'Construct design:', Inches(5.6), size=14, italic=False, color=NAVY)
add_bullets(slide, [
    'CD19 ECD: Met1 to ~Lys291 (UniProt P15391) — remove TM + cytoplasmic',
    'SF-CD19 stabilizing mutations from Laurent 2021 supplementary',
    'Cleavable C-terminal His₁₀ + HRV-3C (or Avi tag for biotinylation)',
    'NEVER run kinetics with Fc still attached — bivalent Fc inflates KD 100-1000× via avidity',
    'SEC-MALS QC every batch — target >99% monomer',
], Inches(6.0), size=12)

add_footer(slide, 'Source: doc 18 §5.2. Reeves 2002 PNAS PMID 12370423 (HEK293S GnTI⁻ origin, verified).')

# =====================================================================
# Slide 10: SPR protocol parameters
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'SPR protocol — built from Seigner 2023 benchmark')
add_subtitle(slide, 'Published FMC63-CD19 KD = 5.1 nM (range 2-6 nM); kon = 1.0 × 10⁵ M⁻¹s⁻¹; koff = 5.3 × 10⁻⁴ s⁻¹', Inches(1.0))

spr_data = [
    ['Parameter', 'Value', 'Rationale'],
    ['Chip', 'Biotin CAPture S Series (Cytiva)', 'Re-loadable; reduces immobilization variability'],
    ['Ligand on chip', 'Biotinylated FMC63-Avi-His scFv (~1000 RU)', 'Site-specific biotinylation; oriented capture'],
    ['Analyte in solution', 'SF-CD19 monomer (5 conc.: 0.5, 4, 20, 100, 500 nM)', 'Single-cycle kinetics (SCK) preferred'],
    ['Buffer', 'PBS + 0.1% BSA + 0.05% Tween-20, pH 7.4', 'Seigner 2023; physiological NaCl'],
    ['Temperature / flow', '25 °C / 30 µL/min', 'Standard'],
    ['Association / dissociation', '600 s / 1200 s', 'koff ≈ 5×10⁻⁴ s⁻¹ → t½ ≈ 22 min, need ≥3 t½'],
    ['Regeneration', '3 M GuHCl + 1 M NaOH, 120 s', 'Validate scFv stability across cycles'],
    ['Fitting', '1:1 Langmuir', 'Confirm with global fit + residuals'],
    ['Replicates', 'n ≥ 3 independent runs per mutant', 'Report mean ± SD'],
]
add_table(slide, spr_data, Inches(1.45), Inches(0.4), SW - Inches(0.8), Inches(4.4), font_size=11)

add_callout_box(slide, 'FMC63 has a 20% diabody dimer equilibrium that SEC alone cannot remove (Seigner 2023). For top 3 mutants run Fab format as orthogonal validation — eliminates the diabody artifact entirely.',
                Inches(6.05), Inches(0.5), SW - Inches(1.0), Inches(0.7),
                fill_color=ORANGE, font_size=12)

add_footer(slide, 'Source: doc 18 §5.3. Seigner et al. 2023 Sci Rep 13:23024 PMID 38155191 (verified).')

# =====================================================================
# Slide 11: Tiered biophysics strategy
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Tiered biophysics strategy — pick the right number of mutants per tier')

tier_data = [
    ['Tier', 'Method', '# mutants', 'Throughput', 'Purpose', 'Where'],
    ['1', 'Functional screens (in vivo + in vitro)', '50-100', '–', 'Identify functionally interesting variants', 'CSIR-IGIB'],
    ['2', 'BLI / Octet (FortéBio)', '20-30', '24 / run', 'Rapid kinetic triage', 'Confirm in-house instrument'],
    ['3', 'SPR (Biacore T200 / 8K)', '6-10', '6-30 / day', 'Primary 3D kinetic dataset', 'CSIR-IGIB / NCBS / IISc / IIT-B'],
    ['4', '2D adhesion frequency (Chesla-Zhu)', '2-3', '5-10 / month', 'Membrane-context kinetics; novelty (no published CAR data)', 'External collab — Zhu (Georgia Tech), Liu (Utah), Qi (Tsinghua)'],
    ['5', 'BFP catch bond (Phase 2)', '1-2', '1-2 / quarter', 'Mechanobiology; secondary novelty', 'Same external collaborators'],
]
add_table(slide, tier_data, Inches(1.2), Inches(0.4), SW - Inches(0.8), Inches(3.6), font_size=11)

add_callout_box(slide, 'Key insight: NO published 2D adhesion frequency or BFP catch-bond data exists for ANY CAR-antigen system. This is genuinely novel for FMC63-CD19 — strengthens novelty argument alongside the S214 + Trp212 residue novelty.',
                Inches(5.0), Inches(0.5), SW - Inches(1.0), Inches(0.85),
                fill_color=NAVY, font_size=13)

# Citations
add_subtitle(slide, 'Origin references (verified):', Inches(5.95), size=14, italic=False, color=NAVY)
add_bullets(slide, [
    '2D adhesion frequency: Chesla, Selvaraj, Zhu 1998 Biophys J 75:1553-72 PMID 9726957',
    'BFP catch bonds (TCR): Liu, Chen, Evavold, Zhu 2014 Cell 157:357-68 PMID 24725404',
    'TCR-pMHC 2D kinetics establishing T cell relevance: Huang et al. 2010 Nature 464:932 PMID 20357766',
], Inches(6.4), size=11)

add_footer(slide, 'Source: doc 18 §5.5, §5.7')

# =====================================================================
# Slide 12: Citation correction (He, not Singh)
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Citation correction — He et al. 2023 (NOT Singh)')
add_subtitle(slide, 'PDB 7URV / Sci Immunol cryo-EM paper has first author Changhao He', Inches(1.0))

# Two-column compare
add_callout_box(slide, 'WRONG (current repo)\n\nSingh et al. 2023, Sci Immunol\n(used in 16 markdown files + slide deck)',
                Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.0),
                fill_color=RED_ACCENT, font_size=14)

add_callout_box(slide, 'CORRECT (verified PubMed)\n\nHe C, Mansilla-Soto J, Khanra N, Hamieh M, Bustos V, Paquette AJ, Garcia Angus A, Shore DM, Rice WJ, Khelashvili G, Sadelain M, Meyerson JR.\n\nCD19 CAR antigen engagement mechanisms and affinity tuning.\n\nSci Immunol 8(81):eadf1426, 2023.\nPMID 36867678 | PMCID PMC10228544\nDOI 10.1126/sciimmunol.adf1426',
                Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.0),
                fill_color=GREEN_ACCENT, font_size=12)

# Other corrections
add_subtitle(slide, 'Other PMIDs verified during this review:', Inches(3.7), size=14, italic=False, color=NAVY, left=Inches(0.5), width=Inches(6.0))
add_bullets(slide, [
    'Nicholson 1997 (FMC63 origin): PMID 9566763 (Mol Immunol 34:1157-65)',
    'Drent 2019: Clin Cancer Res 25:4014-25, PMID 30979735 (already correct in repo)',
    'Reeves 2002 (HEK293S GnTI⁻): PMID 12370423 (PNAS 99:13419)',
    'Chesla, Selvaraj, Zhu 1998 (2D assay): PMID 9726957 (Biophys J 75:1553-72)',
    'Seigner 2023 (FMC63 KD = 5.1 nM): PMID 38155191 (Sci Rep 13:23024) — confirmed',
    'Laurent 2021 (SF-CD19): PMID 33843201 (ACS Synth Biol 10:1184-98) — first author Laurent, not Zajc',
], Inches(4.1), size=11, left=Inches(0.5), width=Inches(6.0))

add_footer(slide, 'Action: find/replace "Singh et al. 2023" → "He et al. 2023" across 16 markdown files + slide deck. PMID, journal, DOI all unchanged.')

# =====================================================================
# Slide 13: Summary recommendations
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Summary — prioritized recommendations')

rec_data = [
    ['Priority', 'Issue', 'Recommended change'],
    ['CRITICAL', 'Citation: "Singh 2023"', 'Correct to He et al. 2023 across all repo documents'],
    ['CRITICAL', 'Sampling timing', 'Day 1, 2, 4 = activation/exhaustion (correct). Memory pool from terminal sacrifice (post-rest)'],
    ['CRITICAL', '"Humanized mouse"', 'Use NSG (or NSG-MHC-DKO if rest >4 wk); not PBMC-humanized'],
    ['CRITICAL', 'Pre-injection library NGS', 'Sequence plasmid + infusion product as reference'],
    ['CRITICAL', 'CD19 "HEK293S"', 'Switch to Expi293F + SF-CD19 (Laurent 2021)'],
    ['HIGH', 'Library tracking', 'DNA-barcode the library (or confirm single-amplicon coverage)'],
    ['HIGH', 'MOI control', 'MOI ≤ 0.3; sort low-MFI singlets; qPCR VCN'],
    ['HIGH', 'CD19-KO NALM-6 control', 'Add as antigen-escape internal control'],
    ['HIGH', 'Mouse n', 'n ≥ 5 per arm; pre-registered mixed-effects model'],
    ['MEDIUM', 'Fab orthogonal SPR', 'Top 3 mutants — eliminates diabody artifact'],
    ['MEDIUM', 'Octet/BLI triage', 'Add before Biacore commitment'],
    ['MEDIUM', '2D / BFP', 'External collaboration; top 2-3 mutants only; defer BFP to Phase 2'],
    ['MEDIUM', 'Pooled in vivo vs in vitro + arrayed in vivo', 'Consider safer alternative if scope is too ambitious'],
]
add_table(slide, rec_data, Inches(0.45), Inches(0.4), SW - Inches(0.9), Inches(5.7), font_size=10)

add_footer(slide, 'Source: doc 18 §7. Full master reference list: doc 18 §9.')

# =====================================================================
# Slide 14: References
# =====================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, 'Verified references (PubMed-confirmed)')

refs = [
    'Brehm MA et al. NSG-MHC-DKO no xeno-GvHD. FASEB J 33:3137-51, 2019. PMID 30383447.',
    'Brentjens RJ et al. NALM-6 SCID-Beige original. Nat Med 9:279-86, 2003. PMID 12579196.',
    'Brentjens RJ et al. FMC63 19-28z preclinical. Clin Cancer Res 13:5426-35, 2007. PMID 17855649.',
    'Chesla SE, Selvaraj P, Zhu C. 2D adhesion frequency assay origin. Biophys J 75:1553-72, 1998. PMID 9726957.',
    'Drent E et al. Combined CD28+4-1BB affinity-tuned CARs. Clin Cancer Res 25:4014-25, 2019. PMID 30979735.',
    'Eyquem J et al. TRAC knock-in CAR-T. Nature 543:113-7, 2017. PMID 28225754.',
    'Gattinoni L et al. Tscm definition (CD95+). Nat Med 17:1290-7, 2011. PMID 21926977.',
    'He C et al. CD19 CAR engagement / PDB 7URV / KD 4.5 nM. Sci Immunol 8:eadf1426, 2023. PMID 36867678. (NOT Singh)',
    'Huang J et al. 2D TCR-pMHC kinetics. Nature 464:932-6, 2010. PMID 20357766.',
    'Kivioja T et al. UMIs. Nat Methods 9:72-4, 2011. PMID 22101854.',
    'Laurent E et al. SF-CD19 stabilized monomer. ACS Synth Biol 10:1184-98, 2021. PMID 33843201.',
    'Li W et al. MAGeCK. Genome Biol 15:554, 2014. PMID 25476604.',
    'Liu B et al. TCR catch bonds (BFP). Cell 157:357-68, 2014. PMID 24725404.',
    'Long AH et al. 4-1BB ameliorates CAR-T tonic exhaustion. Nat Med 21:581-90, 2015. PMID 25939063.',
    'Mackay LK et al. CD8 Trm pathway. Nat Immunol 14:1294-301, 2013. PMID 24162776.',
    'Majzner RG et al. CD19 antigen density threshold. Cancer Discov 10:702-23, 2020. PMID 32193224.',
    'Milone MC et al. 4-1BB CAR + NALM-6 NSG. Mol Ther 17:1453-64, 2009. PMID 19384291.',
    'Nicholson IC et al. Original FMC63 scFv. Mol Immunol 34:1157-65, 1997. PMID 9566763.',
    'Reeves PJ et al. HEK293S GnTI- origin. PNAS 99:13419-24, 2002. PMID 12370423.',
    'Roth TL et al. Non-viral CRISPR T-cell. Nature 559:405-9, 2018. PMID 30022017.',
    'Sabatino M et al. Clinical-grade CD19 CAR Tscm. Blood 128:519-28, 2016. PMID 27226436.',
    'Sallusto F et al. Tcm/Tem original definition. Nature 401:708-12, 1999. PMID 10537110.',
    'Seigner J et al. FMC63-CD19 KD = 5.1 nM. Sci Rep 13:23024, 2023. PMID 38155191.',
    'Shultz LD et al. Humanized mouse review. Nat Rev Immunol 12:786-98, 2012. PMID 23059428.',
    'Sotillo E et al. CD19 antigen escape. Cancer Discov 5:1282-95, 2015. PMID 26583447.',
    'Wunderlich M et al. NSGS reconstitution. PLoS One 13:e0209034, 2018. PMID 30586420.',
    'Xu Y et al. Tscm correlate with CAR-T persistence. Blood 123:3750-9, 2014. PMID 24782509.',
]

# Two-column refs
col_w = (SW - Inches(1.4)) / 2
left_col = Inches(0.5)
right_col = Inches(0.5) + col_w + Inches(0.4)
half = (len(refs) + 1) // 2

box_l = slide.shapes.add_textbox(left_col, Inches(1.0), col_w, SH - Inches(1.5))
tf = box_l.text_frame
tf.word_wrap = True
for i, r in enumerate(refs[:half]):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = u'• ' + r
    for run in p.runs:
        run.font.size = Pt(9)
        run.font.color.rgb = DARK_GRAY
        run.font.name = 'Calibri'
    p.space_after = Pt(2)

box_r = slide.shapes.add_textbox(right_col, Inches(1.0), col_w, SH - Inches(1.5))
tf = box_r.text_frame
tf.word_wrap = True
for i, r in enumerate(refs[half:]):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = u'• ' + r
    for run in p.runs:
        run.font.size = Pt(9)
        run.font.color.rgb = DARK_GRAY
        run.font.name = 'Calibri'
    p.space_after = Pt(2)

add_footer(slide, 'All PMIDs directly verified against PubMed during preparation of doc 18.')

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
