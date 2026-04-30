#!/usr/bin/env python3
"""
Build a SINGLE Specific Aims slide that Manpreet can drop into
DAC_Presentation_v2.pptx as her Outline / Specific Aims slide.

Output: Specific_Aims_Slide.pptx (one slide, 13.33 x 7.5", 16:9 widescreen)

Design:
  - Top navy title band
  - Project title + student/PI subtitle
  - 2x2 grid of 4 aim cards (Aim 4 = optional/translational)
  - Bottom band: central question / executive pitch
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/Specific_Aims_Slide.pptx'

# Color palette (matching her existing decks)
NAVY = RGBColor(0x1A, 0x23, 0x7E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
TEAL = RGBColor(0x16, 0xA0, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ---------------------------------------------------------------------
# Top title band
# ---------------------------------------------------------------------
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.05))
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.10), SW - Inches(1.0), Inches(0.55))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Specific Aims'
for r in p.runs:
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = 'Calibri'

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), SW - Inches(1.0), Inches(0.40))
tf = subtitle_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ('Advancing CAR-T Cell Therapy by Understanding the Kinetics of Ag-Ab '
          'Interaction Parameters  |  Manpreet Kour, PhD scholar, CSIR-IGIB  |  '
          'PI: Dr. Kausik Chakraborty')
for r in p.runs:
    r.font.size = Pt(13)
    r.font.italic = True
    r.font.color.rgb = WHITE
    r.font.name = 'Calibri'


# ---------------------------------------------------------------------
# Helper to draw an aim card
# ---------------------------------------------------------------------
def add_aim_card(left, top, width, height, aim_num, title, body, question_link, accent_color, optional=False):
    # Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    # Color stripe on left
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.2), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    # Aim number — large
    num_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.1), Inches(1.3), Inches(0.55))
    tf = num_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    if optional:
        p.text = f'AIM {aim_num} *'
    else:
        p.text = f'AIM {aim_num}'
    for r in p.runs:
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = accent_color
        r.font.name = 'Calibri'

    # Question link tag (top right of card)
    tag_box = slide.shapes.add_textbox(left + width - Inches(2.0), top + Inches(0.1), Inches(1.85), Inches(0.35))
    tf = tag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question_link
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(10)
        r.font.italic = True
        r.font.color.rgb = DARK_GRAY
        r.font.name = 'Calibri'

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.65), width - Inches(0.5), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    for r in p.runs:
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = NAVY
        r.font.name = 'Calibri'

    # Body
    body_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.20), width - Inches(0.5), height - Inches(1.30))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if line.startswith('•'):
            p.text = line
        else:
            p.text = line
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.color.rgb = DARK_GRAY
            r.font.name = 'Calibri'
        p.space_after = Pt(2)


# ---------------------------------------------------------------------
# 2x2 grid of aim cards
# ---------------------------------------------------------------------
# Card layout: ~6.4" wide x ~2.5" tall, with ~0.2" gap
card_w = Inches(6.35)
card_h = Inches(2.45)
gap = Inches(0.18)
left_x = Inches(0.30)
right_x = left_x + card_w + gap
top_y = Inches(1.30)
bot_y = top_y + card_h + gap

# Aim 1 — Library design + computational
add_aim_card(
    left_x, top_y, card_w, card_h,
    aim_num='1',
    title='Library design and computational validation',
    body=[
        '• Identify and computationally validate critical scFv contact residues at the FMC63-CD19 interface (PDB 7URV; He et al. 2023, PMID 36867678) using structural analysis (FreeSASA, contact mapping) and mCSM-AB2 ΔΔG predictions.',
        '• Design NNK saturation library at 4 positions:',
        '    PRIMARY (novel — never published): S214, Trp212',
        '    SECONDARY (beyond He et al.\'s single Ala): Y260, Y261',
        '• Library size: ~376 single-position variants spanning the full affinity spectrum.',
    ],
    question_link='[Foundation — already done]',
    accent_color=PURPLE,
)

# Aim 2 — Functional mapping
add_aim_card(
    right_x, top_y, card_w, card_h,
    aim_num='2',
    title='Functional mapping in vitro and in vivo',
    body=[
        '• Map CAR-T activation kinetics (CD69, CD25, IFN-γ at 24/48/96 h), exhaustion trajectory (PD-1, TIM-3, LAG-3, TOX), cytotoxic function, and memory formation across the affinity variant library.',
        '• Models:  in vitro NALM-6 co-culture with graded CD19 density (CD19-KO + titrated re-expression; Majzner 2020 PMID 32193224)',
        '              in vivo NALM-6/NSG xenograft (or NSG-MHC-DKO if rest >4 wk; Brehm 2019 PMID 30383447)',
        '• Readout: sort-then-NGS to identify variants enriched in each phenotype (MAGeCK MLE).',
    ],
    question_link='[Q1 reformulated]',
    accent_color=GREEN_ACCENT,
)

# Aim 3 — Biophysical characterization
add_aim_card(
    left_x, bot_y, card_w, card_h,
    aim_num='3',
    title='Biophysical characterization and kinetic-functional correlation',
    body=[
        '• Kinetically characterize a panel of variants spanning the full functional spectrum (high / intermediate / low / WT controls).',
        '• Methods:  SPR (Biacore) — 3D KD, kon, koff, dwell time (benchmark Seigner 2023, PMID 38155191: KD = 5.1 nM)',
        '              2D micropipette adhesion frequency (Chesla-Zhu 1998 PMID 9726957) — membrane-context kinetics; novel for any CAR-antigen system',
        '• Statistically correlate kinetic parameters with functional readouts to define the optimal kinetic window.',
    ],
    question_link='[Q2 reformulated]',
    accent_color=TEAL,
)

# Aim 4 — Translational (optional)
add_aim_card(
    right_x, bot_y, card_w, card_h,
    aim_num='4',
    title='Translational validation (optional / stretch)',
    body=[
        '• Validate top-performing affinity variants in primary human T cells (multiple allogeneic donors).',
        '• Assess in vivo persistence and memory recall in NSG-MHC-DKO xenograft models (Brehm 2019).',
        '• Translational endpoints: tumor-free survival, CAR-T peripheral persistence, memory subset durability post-rechallenge.',
        '* Optional aim — pursued if scope and timeline allow.',
    ],
    question_link='[Stretch goal]',
    accent_color=ORANGE,
    optional=True,
)


# ---------------------------------------------------------------------
# Bottom executive pitch band
# ---------------------------------------------------------------------
pitch_top = Inches(6.55)
pitch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), pitch_top, SW - Inches(0.6), Inches(0.70))
pitch.fill.solid()
pitch.fill.fore_color.rgb = NAVY
pitch.line.fill.background()

pitch_tf = pitch.text_frame
pitch_tf.word_wrap = True
pitch_tf.margin_left = Inches(0.20)
pitch_tf.margin_right = Inches(0.20)
pitch_tf.margin_top = Inches(0.05)
pitch_tf.margin_bottom = Inches(0.05)
p = pitch_tf.paragraphs[0]
p.text = ('Central contribution:  the first systematic dataset linking ~376 FMC63 scFv variants '
          'to a full functional matrix (activation / exhaustion / memory / persistence) and a '
          'multi-parameter kinetic dataset (3D + 2D), defining the kinetic-function relationship '
          'for anti-CD19 CAR-T therapy.')
p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(12)
    r.font.italic = True
    r.font.color.rgb = WHITE
    r.font.name = 'Calibri'


# ---------------------------------------------------------------------
# Footer
# ---------------------------------------------------------------------
footer = slide.shapes.add_textbox(Inches(0.30), SH - Inches(0.30), SW - Inches(0.6), Inches(0.25))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = 'Two of the four mutated positions (S214, Trp212) have never been mutated in any published CAR study  •  2D adhesion frequency / catch-bond data have never been published for any CAR-antigen system'
p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)
    r.font.name = 'Calibri'


prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width / Inches(1):.3f} x {prs.slide_height / Inches(1):.3f} inches")
