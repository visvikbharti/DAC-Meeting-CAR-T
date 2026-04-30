#!/usr/bin/env python3
"""
Update Manpreet's DAC1_draft2.pptx with her stated objectives.

Slide 20 (Objectives):
  - Slot 1 (oval '1') — populate with Objective 1 (currently empty)
  - Slot 2 (oval '2') — populate with Objective 2 (currently empty)
  - Slot 3 (oval '3') — leave existing text untouched ("To Understand the
    Parameters of Ag-Ab Binding which Affects the Amplitude of the Signal")

Slide 27 — heading currently duplicates slide 21-26 ("OBJECTIVE 1.2 To
develop scFv mutant library..."), but content is the in vivo experimental
design. Fix the heading to match the actual content.

Backup at DAC1_draft2.backup.pptx (already created).
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree

PPTX = '/Users/vishalbharti/Downloads/DAC-Meeting-CAR-T-MANPREET/DAC1_draft2.pptx'

OBJ1_TEXT = '    To determine CAR-T cell efficacy in vivo and in vitro using scFv mutant CAR library'
OBJ2_TEXT = '    To do biophysical characterization of scFv mutants which performed better in vivo as well as in vitro'

# Slide 27 heading correction
OLD_HEADING_27 = 'OBJECTIVE 1.2 To develop scFv mutant library of CAR & make mutant CARPOOL T cells'
NEW_HEADING_27 = 'OBJECTIVE 1.2 To test mutant CARPOOL T cells in vivo using luciferinized NALM-6 mouse model'


def set_rectangle_text(shape, text):
    """Set text on an empty rectangle shape, matching slot-3 format
    (Times New Roman, 22pt, bold, vertically centered)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Clear existing paragraphs (if any)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = 'Times New Roman'
    run.font.size = Pt(22)
    run.font.bold = True


def main():
    p = Presentation(PPTX)

    # ---- Slide 20: Objectives slide ----
    slide20 = p.slides[19]
    # Shape 0 = Rectangle 2 (slot 1, currently empty)
    # Shape 2 = Rectangle 4 (slot 2, currently empty)
    # Shape 4 = Rectangle 6 (slot 3, has existing text)
    rect_slot1 = slide20.shapes[0]
    rect_slot2 = slide20.shapes[2]

    # Sanity check shape names
    assert rect_slot1.name == 'Rectangle 2', f'unexpected shape: {rect_slot1.name}'
    assert rect_slot2.name == 'Rectangle 4', f'unexpected shape: {rect_slot2.name}'

    set_rectangle_text(rect_slot1, OBJ1_TEXT)
    set_rectangle_text(rect_slot2, OBJ2_TEXT)
    print(f'  ✔ Slide 20 slot 1 populated: "{OBJ1_TEXT.strip()}"')
    print(f'  ✔ Slide 20 slot 2 populated: "{OBJ2_TEXT.strip()}"')

    # ---- Slide 27: heading correction ----
    slide27 = p.slides[26]
    fixed = False
    for sh in slide27.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == OLD_HEADING_27:
                        run.text = NEW_HEADING_27
                        fixed = True
                        print(f'  ✔ Slide 27 heading corrected')
                        break
                if fixed:
                    break
        if fixed:
            break
    if not fixed:
        # Heading text might span multiple runs; try whole-paragraph search
        for sh in slide27.shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    full = ''.join(r.text for r in para.runs)
                    if 'OBJECTIVE 1.2 To develop scFv mutant library' in full:
                        # Replace first run, clear others
                        if para.runs:
                            para.runs[0].text = NEW_HEADING_27
                            for r in para.runs[1:]:
                                r.text = ''
                            fixed = True
                            print(f'  ✔ Slide 27 heading corrected (multi-run)')
                            break
                if fixed:
                    break
            if fixed:
                break
    if not fixed:
        print('  ⚠ Slide 27 heading not found — leaving unchanged')

    p.save(PPTX)
    print(f'\nSaved: {PPTX}')


if __name__ == '__main__':
    main()
